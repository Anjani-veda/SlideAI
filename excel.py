import os
import uuid
import re
import pandas as pd
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from docx import Document
from docx.shared import Inches as DocxInches

from fpdf import FPDF


# ---------------- PATH SETUP ---------------- #
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

REPORT_DIR = os.path.join(BASE_DIR, "reports")
STATIC_DIR = os.path.join(BASE_DIR, "static")

os.makedirs(REPORT_DIR, exist_ok=True)
os.makedirs(STATIC_DIR, exist_ok=True)


# ---------------- THEMES ---------------- #
THEMES = {
    "Emerald": {
        "p": (33, 115, 70),
        "bg": (220, 252, 231),
        "d": (20, 83, 45)
    },

    "Light Purple": {
        "p": (109, 40, 217),
        "bg": (245, 243, 255),
        "d": (76, 29, 149)
    },

    "Dark Blue": {
        "p": (30, 58, 138),
        "bg": (219, 234, 254),
        "d": (23, 37, 84)
    },

    "Gradient Blue": {
        "p": (2, 132, 199),
        "bg": (224, 242, 254),
        "d": (12, 74, 110)
    },

    "Midnight": {
        "p": (15, 23, 42),
        "bg": (226, 232, 240),
        "d": (2, 6, 23)
    },

    "Sunset": {
        "p": (194, 65, 12),
        "bg": (255, 247, 237),
        "d": (124, 45, 18)
    }
}


# ---------------- UTILS ---------------- #
def strip_emojis(text):
    return re.sub(r'[^\x00-\x7F]+', '', str(text)).strip()


def safe_filename(name):
    clean = re.sub(r'[\\/*?:"<>|]', "", str(name))
    return clean[:50] if clean else "data"


def load_file(path):
    ext = os.path.splitext(path)[1].lower()

    if ext == '.csv':
        return {"Sheet1": pd.read_csv(path)}

    return pd.read_excel(path, sheet_name=None)


def clean_data(df):

    if df is None or df.empty:
        return pd.DataFrame()

    df = df.dropna(how='all')
    df = df.dropna(axis=1, how='all')
    df = df.reset_index(drop=True)

    new_cols = []
    seen = set()

    for i, col in enumerate(df.columns):

        base = str(col).strip()

        if "Unnamed" in base or not base:
            base = f"Column_{i}"

        final = base
        c = 1

        while final in seen:
            final = f"{base}_{c}"
            c += 1

        seen.add(final)
        new_cols.append(final)

    df.columns = new_cols

    return df


def get_column_types(df):

    num = []
    text = []

    for col in df.columns:

        if any(
            k in str(col).lower()
            for k in ["date", "year", "month", "no", "id", "code", "sl"]
        ):
            text.append(col)
            continue

        if not pd.api.types.is_numeric_dtype(df[col]):

            try:
                clean_s = df[col].astype(str).str.replace(
                    r'[$,₹,£,€,%, ]',
                    '',
                    regex=True
                )

                temp_num = pd.to_numeric(clean_s, errors='coerce')

                if temp_num.notna().mean() > 0.5:
                    df[col] = temp_num

            except:
                pass

        if pd.api.types.is_numeric_dtype(df[col]):
            num.append(col)
        else:
            text.append(col)

    return num, text


# ---------------- ANALYSIS ---------------- #
def generate_kpis(df, num_cols, txt_cols):

    kpis = {}

    if df.empty:
        return kpis

    def _most_repeated(col):

        clean_s = df[col].astype(str).str.strip()

        clean_s = clean_s[
            ~clean_s.str.lower().isin(["nan", "none", ""])
        ]

        if clean_s.empty:
            return None, 0

        def get_pattern(x):

            s = re.sub(r'[^a-zA-Z0-9\s]', ' ', str(x))
            words = s.split()

            return " ".join(words[:3]).lower() if words else ""

        patterns = clean_s.apply(get_pattern)
        patterns = patterns[patterns != ""]

        if patterns.empty:
            return None, 0

        top_counts = patterns.value_counts()

        if top_counts.empty:
            return None, 0

        top = top_counts.index[0]
        cnt = int(top_counts.iloc[0])

        matches = clean_s[patterns == top].value_counts()

        if matches.empty:
            return None, 0

        orig = matches.index[0]

        return orig, cnt

    target = next(
        (
            c for c in txt_cols
            if "name" in c.lower() or "client" in c.lower()
        ),
        txt_cols[0] if txt_cols else None
    )

    if target:

        orig, count = _most_repeated(target)

        if orig:
            kpis[f"Most Repeated ({target})"] = {
                "Current": f"{orig} ({count} times)",
                "Trend": "Frequency"
            }

    for col in num_cols:

        c_low = str(col).lower()

        if any(
            k in c_low
            for k in [
                "invoice",
                "value",
                "amount",
                "total",
                "price",
                "revenue",
                "cost",
                "amt"
            ]
        ):

            try:
                v_numeric = pd.to_numeric(df[col], errors='coerce')

                v_total = v_numeric.sum()
                v_max = v_numeric.max()
                v_min = v_numeric.min()

                if pd.notna(v_total):
                    kpis[f"Total {col}"] = {
                        "Current": f"{v_total:,.2f}",
                        "Trend": f"Sum of {col}"
                    }

                if pd.notna(v_max):
                    kpis[f"Max {col}"] = {
                        "Current": f"{v_max:,.2f}",
                        "Trend": f"Highest in {col}"
                    }

                if pd.notna(v_min):
                    kpis[f"Min {col}"] = {
                        "Current": f"{v_min:,.2f}",
                        "Trend": f"Lowest in {col}"
                    }

            except:
                pass

    return kpis


def get_overall_stats(df, num_cols):

    stats = {
        "Total Value": 0.0,
        "Average": 0.0,
        "Max Value": 0.0,
        "Min Value": 0.0,
        "Records": len(df)
    }

    if df.empty or not num_cols:
        return stats

    target = None

    for col in num_cols:

        c_low = str(col).lower()

        if any(
            k in c_low
            for k in [
                "total",
                "amount",
                "revenue",
                "price",
                "value",
                "amt"
            ]
        ):
            target = col
            break

    if not target:
        target = num_cols[0]

    try:
        v_numeric = pd.to_numeric(df[target], errors='coerce')

        stats["Total Value"] = float(v_numeric.sum())
        stats["Average"] = float(v_numeric.mean())
        stats["Max Value"] = float(v_numeric.max())
        stats["Min Value"] = float(v_numeric.min())

    except:
        pass

    return stats


# ---------------- CHARTS ---------------- #
def charts(df, num_cols, txt_cols, rid):

    path = os.path.join(STATIC_DIR, str(rid))
    os.makedirs(path, exist_ok=True)

    res = {}

    if df.empty:
        return res

    # BAR + PIE
    if txt_cols:

        target_col = txt_cols[0]
        counts = df[target_col].value_counts().head(5)

        if not counts.empty:

            # BAR
            plt.figure(figsize=(8, 5))

            counts.plot(
                kind='bar',
                color='#217346'
            )

            plt.title(f"Top 5: {target_col}")
            plt.tight_layout()

            f_bar = os.path.join(path, "bar.png")

            plt.savefig(f_bar)
            plt.close()

            res["Bar Chart"] = {
                "path": f_bar,
                "conclusion": "Top categories identified."
            }

            # PIE
            plt.figure(figsize=(8, 5))

            counts.plot(
                kind='pie',
                autopct='%1.1f%%'
            )

            plt.title(f"Distribution: {target_col}")
            plt.ylabel("")

            plt.tight_layout()

            f_pie = os.path.join(path, "pie.png")

            plt.savefig(f_pie)
            plt.close()

            res["Pie Chart"] = {
                "path": f_pie,
                "conclusion": "Distribution visualized."
            }

    # LINE
    if num_cols:

        target_num = num_cols[0]
        subset = df[target_num].head(20)

        if not subset.empty:

            plt.figure(figsize=(8, 5))

            subset.plot(
                kind='line',
                marker='o',
                color='#217346'
            )

            plt.title(f"Trend: {target_num}")

            plt.tight_layout()

            f_line = os.path.join(path, "line.png")

            plt.savefig(f_line)
            plt.close()

            res["Line Chart"] = {
                "path": f_line,
                "conclusion": "Trend identified."
            }

    return res


# ---------------- PPT REPORT ---------------- #
def create_ppt(
    kpis,
    charts_data,
    theme_name,
    client_data,
    top_table=None,
    insights=None,
    font_name="Arial"
):

    os.makedirs(REPORT_DIR, exist_ok=True)

    prs = Presentation()

    t = THEMES.get(theme_name, THEMES["Emerald"])

    PRI = RGBColor(*t["p"])
    BG = RGBColor(*t["bg"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box = slide.shapes.add_shape(
        1,
        0,
        0,
        prs.slide_width,
        prs.slide_height
    )

    box.fill.solid()
    box.fill.fore_color.rgb = PRI

    # KPI SLIDE
    if kpis:

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        for i, (l, d) in enumerate(list(kpis.items())[:8]):

            x = Inches(0.3 + (i % 2) * 4.8)
            y = Inches(1 + (i // 2) * 1.5)

            b = slide.shapes.add_shape(
                1,
                x,
                y,
                Inches(4.5),
                Inches(1.2)
            )

            b.fill.solid()
            b.fill.fore_color.rgb = BG

            tx = slide.shapes.add_textbox(
                x,
                y,
                Inches(4.5),
                Inches(1.2)
            )

            p = tx.text_frame.paragraphs[0]

            p.text = f"{strip_emojis(l)}: {d['Current']}"

            p.font.size = Pt(14)
            p.font.name = font_name

    # CHART SLIDES
    for k, d in charts_data.items():

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        slide.shapes.add_picture(
            d['path'],
            Inches(1),
            Inches(1),
            width=Inches(11)
        )

    f = os.path.join(
        REPORT_DIR,
        f"report_{uuid.uuid4()}.pptx"
    )

    prs.save(f)

    return f


# ---------------- DOCX REPORT ---------------- #
def create_docx(
    kpis,
    charts_data,
    theme_name,
    client_data,
    top_table=None,
    insights=None,
    font_name="Arial"
):

    os.makedirs(REPORT_DIR, exist_ok=True)

    doc = Document()

    doc.add_heading('Analysis Report', 0)

    for l, d in kpis.items():
        doc.add_paragraph(
            f"{strip_emojis(l)}: {d['Current']}"
        )

    for k, d in charts_data.items():

        doc.add_heading(k, 1)

        doc.add_picture(
            d['path'],
            width=DocxInches(5)
        )

    f = os.path.join(
        REPORT_DIR,
        f"report_{uuid.uuid4()}.docx"
    )

    doc.save(f)

    return f


# ---------------- PDF REPORT ---------------- #
def create_pdf(
    kpis,
    charts_data,
    theme_name,
    client_data,
    top_table=None,
    insights=None,
    font_name="Arial"
):

    os.makedirs(REPORT_DIR, exist_ok=True)

    pdf = FPDF()

    pdf.add_page()

    pdf.set_font("Helvetica", 'B', 16)

    pdf.cell(
        190,
        10,
        text="Analysis Report",
        align='C'
    )

    pdf.ln(10)

    pdf.set_font("Helvetica", '', 10)

    for l, d in kpis.items():

        pdf.multi_cell(
            190,
            7,
            text=f"{strip_emojis(l)}: {d['Current']}",
            border=1
        )

    for k, d in charts_data.items():

        pdf.add_page()

        pdf.image(
            d['path'],
            x=10,
            w=180
        )

    f = os.path.join(
        REPORT_DIR,
        f"report_{uuid.uuid4()}.pdf"
    )

    pdf.output(f)

    return f


# ---------------- PROCESS ---------------- #
def process_file(path, theme, font=None, options=None):

    try:

        sheets = load_file(path)

        cleaned = {}

        for name, df in sheets.items():

            dfc = clean_data(df)

            if not dfc.empty:
                dfc["Sheet Source"] = name
                cleaned[name] = dfc

        if not cleaned:
            return {
                "status": "error",
                "message": "Empty file."
            }

        order = {}

        if len(cleaned) > 1:

            order["Overall Summary"] = pd.concat(
                cleaned.values(),
                ignore_index=True,
                sort=False
            )

        order.update(cleaned)

        results = {}

        for name, df in order.items():

            num, txt = get_column_types(df)

            kpis = generate_kpis(df, num, txt)

            ch = charts(
                df,
                num,
                txt,
                uuid.uuid4()
            )

            overall_stats = get_overall_stats(df, num)

            results[name] = {

                "kpis": kpis,

                "ppt_path": create_ppt(
                    kpis,
                    ch,
                    theme,
                    [],
                    font_name=font
                ),

                "docx_path": create_docx(
                    kpis,
                    ch,
                    theme,
                    [],
                    font_name=font
                ),

                "pdf_path": create_pdf(
                    kpis,
                    ch,
                    theme,
                    [],
                    font_name=font
                ),

                "overall_metrics": overall_stats,

                "charts": ch,

                "full_data": df.head(500).to_dict('records')
            }

        return {
            "status": "success",
            "sheets": results
        }

    except Exception as e:

        return {
            "status": "error",
            "message": str(e)
        }
